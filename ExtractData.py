from pptx import Presentation

def extract_data_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    summary_data = {}

    # Extract content from each slide
    for i, slide in enumerate(presentation.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)

        # Example: Get specific slides
        if i == 1:  # Slide for Date and Client Name
            summary_data['ClientName'] = " ".join(slide_text)
        elif i == 2:  # Financial Well Being (slide 12)
            summary_data['FinancialWellBeing'] = " ".join(slide_text)
        elif i == 3:  # Key Wins (slide 3)
            summary_data['KeyWins'] = " ".join(slide_text)
        # Add other slides and sections as needed
    
    return summary_data

# Test with a PowerPoint file
pptx_path = 'path_to_your_pptx_file.pptx'
summary_data = extract_data_from_pptx(pptx_path)
print(summary_data)