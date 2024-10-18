import os
from pptx import Presentation
from fpdf import FPDF

def extract_data_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    summary_data = {}

    # Extract content from each slide
    for i, slide in enumerate(presentation.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)

        # Clean and join the slide text
        joined_text = " ".join(slide_text).strip()

        # Get specific slides based on their index
        if i == 0:  # Slide for Date and Client Name (page 1)
            summary_data['ClientName'] = joined_text
        elif i == 1:  # Financial Well Being (page 2)
            summary_data['FinancialWellBeing'] = joined_text
        elif i == 2:  # Key Wins (page 3)
            summary_data['KeyWins'] = joined_text
        elif i == 3:  # Goals (page 4)
            summary_data['Goals'] = joined_text
        elif i == 4:  # Estate Planning (page 5)
            summary_data['EstatePlanning'] = joined_text
        elif i == 5:  # Your Tax Rates (page 6)
            summary_data['YourTaxRates'] = joined_text
        elif i == 6:  # Net Worth (page 7)
            summary_data['NetWorth'] = joined_text
        elif i == 7:  # Asset Allocation (page 8)
            summary_data['AssetAllocation'] = joined_text
        elif i == 8:  # Investment Tax Allocation (page 9)
            summary_data['InvestmentTaxAllocation'] = joined_text
        elif i == 9:  # Homework (page 10)
            summary_data['Homework'] = joined_text

    return summary_data